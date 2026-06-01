"""
Academic Blue — default slide theme for gen-slides.

Provides helper functions and constants for generating clean,
professional academic presentations via python-pptx.
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import platform
import os

# --- Color Palette ---
PRIMARY = RGBColor(0x1F, 0x4E, 0x79)
ACCENT = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BG = RGBColor(0xF2, 0xF7, 0xFC)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)

# --- Typography ---
def get_system_font():
    """Return a safe sans-serif font for the current OS."""
    system = platform.system()
    if system == "Darwin":
        return "PingFang SC"
    elif system == "Windows":
        return "Microsoft YaHei"
    return "Noto Sans CJK SC"

FONT_FAMILY = get_system_font()
TITLE_SIZE = Pt(28)
SUBTITLE_SIZE = Pt(18)
BODY_SIZE = Pt(16)
CAPTION_SIZE = Pt(12)

# --- Layout Constants ---
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN_LEFT = Inches(0.8)
MARGIN_TOP = Inches(1.4)
CONTENT_WIDTH = Inches(11.7)
CONTENT_HEIGHT = Inches(5.5)


def add_title_bar(slide, title_text):
    """Add a colored title bar at the top of a content slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0), top=Inches(0),
        width=SLIDE_WIDTH, height=Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(
        left=MARGIN_LEFT, top=Inches(0.15),
        width=CONTENT_WIDTH, height=Inches(0.8)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = TITLE_SIZE
    run.font.bold = True
    run.font.color.rgb = LIGHT_TEXT
    run.font.name = FONT_FAMILY


def add_logo(slide, logo_path, size=Inches(0.8)):
    """Place a logo at the top-right corner of a slide."""
    if not os.path.exists(logo_path):
        return
    slide.shapes.add_picture(
        logo_path,
        left=SLIDE_WIDTH - size - Inches(0.3),
        top=Inches(0.15),
        height=size
    )


def add_bullets(slide, bullets, top=None, left=None, width=None):
    """Add a bulleted text box to a slide."""
    _top = top or MARGIN_TOP
    _left = left or MARGIN_LEFT
    _width = width or CONTENT_WIDTH

    txBox = slide.shapes.add_textbox(
        left=_left, top=_top,
        width=_width, height=CONTENT_HEIGHT
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.size = BODY_SIZE
        run.font.color.rgb = DARK_TEXT
        run.font.name = FONT_FAMILY


def add_centered_image(slide, img_path, max_width=Inches(9), top=Inches(1.5)):
    """Add a centered image to a slide, scaled to fit."""
    if not os.path.exists(img_path):
        return
    from PIL import Image
    img = Image.open(img_path)
    w, h = img.size
    aspect = w / h

    width = min(max_width, SLIDE_WIDTH - Inches(1.6))
    height = width / aspect
    if height > Inches(5.5):
        height = Inches(5.5)
        width = height * aspect

    left = (SLIDE_WIDTH - width) / 2
    slide.shapes.add_picture(img_path, left=int(left), top=int(top), width=int(width))
